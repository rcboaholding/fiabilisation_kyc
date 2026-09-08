# -*- coding: utf-8 -*-
"""Charte graphique reprise de « exemple rapport.pptx » (Conformité BOA Group).

Grammaire du modèle, respectée à l'identique :
  - format 10 x 5,625 pouces ;
  - couverture : fond vert 009A56, deux ellipses 00C060 à droite ;
  - page courante : fond blanc, deux ellipses C8E6C9, bandeau de titre vert
    plein largeur (hauteur 0,62"), intitulé de bloc vert + filet, pied de page
    gris E0E0E0, logo en haut à droite ;
  - typographie Calibri, corps 11 pt, notes 8–9 pt.
"""
import os

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

GREEN = RGBColor(0x00, 0x9A, 0x56)
GREEN_BRIGHT = RGBColor(0x00, 0xC0, 0x60)
GREEN_PALE = RGBColor(0xC8, 0xE6, 0xC9)
GREEN_SUB = RGBColor(0xC8, 0xF0, 0xD8)
GREEN_DATE = RGBColor(0xD4, 0xF5, 0xE4)
GREEN_FOOT = RGBColor(0xA0, 0xD4, 0xB8)
GREEN_SOFT = RGBColor(0xF0, 0xF7, 0xF4)
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
BODY = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x66, 0x66, 0x66)
FOOT = RGBColor(0x88, 0x88, 0x88)
CARD = RGBColor(0xF5, 0xF5, 0xF5)
DIVIDER = RGBColor(0xE0, 0xE0, 0xE0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_PALE = RGBColor(0xFD, 0xF3, 0xDF)
RED = RGBColor(0xD3, 0x2F, 0x2F)
RED_PALE = RGBColor(0xFB, 0xE9, 0xE7)
CYAN = RGBColor(0x00, 0xB0, 0xF0)
GREEN_OK = RGBColor(0x70, 0xAD, 0x47)

FONT = "Calibri"
W, H = 10.0, 5.625

LEFT = 0.4
RIGHT = 9.6
BAND_H = 0.62
RULE_Y = 1.08
BODY_TOP = 1.22
BODY_BOTTOM = 5.30
FOOTER_Y = 5.35

LOGO = next((p for p in ("assets_tpl/tpl_image1.png",) if os.path.exists(p)), None)


def _set_bg(slide, color):
    """Fond de diapositive plein (le modèle utilise <p:bg> et non une forme)."""
    xml = slide._element
    bg = xml.makeelement(qn("p:bg"), {})
    bgpr = xml.makeelement(qn("p:bgPr"), {})
    fill = xml.makeelement(qn("a:solidFill"), {})
    clr = xml.makeelement(qn("a:srgbClr"), {"val": str(color)})
    fill.append(clr)
    bgpr.append(fill)
    bgpr.append(xml.makeelement(qn("a:effectLst"), {}))
    bg.append(bgpr)
    xml.find(qn("p:cSld")).insert(0, bg)


def rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE,
         line_w=0.75):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    sp.text_frame.word_wrap = True
    return sp


def txt(slide, x, y, w, h, contenu, size=11, bold=False, color=BODY,
        align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP,
        spacing=None, wrap=True, line_spacing=None):
    """contenu : str, liste de str, ou liste de tuples (texte, size, bold, color)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    lignes = contenu if isinstance(contenu, (list, tuple)) else [contenu]
    for i, ligne in enumerate(lignes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.space_after = Pt(spacing)
        if line_spacing:
            p.line_spacing = line_spacing
        if isinstance(ligne, tuple):
            s, sz, b, col = (ligne + (None, None, None))[:4]
        else:
            s, sz, b, col = ligne, None, None, None
        r = p.add_run()
        r.text = s
        r.font.size = Pt(sz or size)
        r.font.bold = bold if b is None else b
        r.font.color.rgb = col or color
        r.font.name = FONT
        r.font.italic = italic
    return tb


def _decor(slide, couleur):
    """Les deux ellipses d'angle du modèle, posées avant tout contenu."""
    rect(slide, 8.2, -1.5, 4.0, 5.0, fill=couleur, shape=MSO_SHAPE.OVAL)
    rect(slide, 8.8, 2.0, 3.0, 3.8, fill=couleur, shape=MSO_SHAPE.OVAL)


def _logo(slide):
    if LOGO:
        slide.shapes.add_picture(LOGO, Inches(7.45), Inches(0.0), Inches(2.55))


def couverture(prs, titre, sous_titre, date_txt, pied):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, GREEN)
    _decor(s, GREEN_BRIGHT)
    _logo(s)
    txt(s, 0.5, 0.85, 7.3, 1.55, titre, size=32, bold=True, color=WHITE,
        anchor=MSO_ANCHOR.MIDDLE, spacing=2)
    txt(s, 0.5, 2.55, 7.5, 0.5, sous_titre, size=14, bold=True, color=GREEN_SUB)
    rect(s, 0.5, 3.18, 5.5, 0.06, fill=WHITE)
    txt(s, 0.5, 3.42, 7.0, 0.35, date_txt, size=11, color=GREEN_DATE)
    txt(s, 0.5, 5.2, 7.0, 0.28, pied, size=8, color=GREEN_FOOT)
    return s


def cloture(prs, titre, sous_titre, date_txt):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, GREEN)
    _decor(s, GREEN_BRIGHT)
    _logo(s)
    txt(s, 1.0, 1.5, 6.0, 1.4, titre, size=54, bold=True, color=WHITE,
        anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 1.0, 3.1, 5.5, 0.06, fill=WHITE)
    txt(s, 1.0, 3.3, 7.5, 0.4, sous_titre, size=13, bold=True, color=GREEN_SUB)
    txt(s, 1.0, 3.85, 7.0, 0.3, date_txt, size=10, color=GREEN_FOOT)
    return s


def page(prs, titre_bandeau, intitule=None, pied=""):
    """Page courante : bandeau vert, intitulé de bloc + filet, pied de page."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, WHITE)
    _decor(s, GREEN_PALE)
    rect(s, 0, 0, W, BAND_H, fill=GREEN)
    txt(s, 0.3, 0, 7.0, BAND_H, titre_bandeau, size=18, bold=True, color=WHITE,
        anchor=MSO_ANCHOR.MIDDLE)
    _logo(s)
    if intitule:
        txt(s, LEFT, 0.78, 8.0, 0.32, intitule.upper(), size=11, bold=True, color=GREEN)
        rect(s, LEFT, RULE_Y, 9.2, 0.04, fill=GREEN)
    rect(s, 0, FOOTER_Y, W, 0.28, fill=DIVIDER)
    txt(s, 6.0, FOOTER_Y, 3.4, 0.28, pied, size=8, color=FOOT,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    return s


def sommaire(prs, entrees, pied=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, WHITE)
    _decor(s, GREEN_PALE)
    _logo(s)
    txt(s, 0.5, 0.5, 5.5, 1.1, "SOMMAIRE", size=44, bold=True, color=GREEN,
        anchor=MSO_ANCHOR.MIDDLE)
    y = 1.75
    for rom, libelle in entrees:
        rect(s, 0.5, y + 0.08, 0.07, 0.38, fill=GREEN)
        txt(s, 0.65, y, 0.6, 0.55, rom, size=11, bold=True, color=GREEN,
            anchor=MSO_ANCHOR.MIDDLE)
        txt(s, 1.3, y, 7.0, 0.55, libelle, size=14, bold=True, color=NAVY,
            anchor=MSO_ANCHOR.MIDDLE)
        rect(s, 0.5, y + 0.50, 6.5, 0.03, fill=DIVIDER)
        y += 0.62
    rect(s, 0, FOOTER_Y, W, 0.28, fill=DIVIDER)
    txt(s, 6.0, FOOTER_Y, 3.4, 0.28, pied, size=8, color=FOOT,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    return s


def bloc(slide, y, intitule):
    """Second intitulé de bloc (vert + filet) à l'intérieur d'une page."""
    txt(slide, LEFT, y, 8.0, 0.3, intitule.upper(), size=11, bold=True, color=GREEN)
    rect(slide, LEFT, y + 0.28, 9.2, 0.04, fill=GREEN)
    return y + 0.42


KPI_H = 0.86


def kpi(slide, x, y, w, h=KPI_H, valeur="", libelle="", note=None, couleur=GREEN,
        vsize=20):
    """Carte chiffre-clé : bandeau de couleur, valeur en haut, libellé et note
    ancrés au bas de la carte — le texte ne peut donc pas déborder."""
    rect(slide, x, y, w, h, fill=CARD)
    rect(slide, x, y, w, 0.05, fill=couleur)
    txt(slide, x + 0.12, y + 0.07, w - 0.24, vsize * 1.30 / 72.0, valeur,
        size=vsize, bold=True, color=couleur)
    if note:
        txt(slide, x + 0.12, y + h - 0.38, w - 0.24, 0.18, libelle,
            size=8.5, bold=True, color=NAVY)
        txt(slide, x + 0.12, y + h - 0.19, w - 0.24, 0.17, note,
            size=7, color=MUTED)
    else:
        txt(slide, x + 0.12, y + h - 0.26, w - 0.24, 0.20, libelle,
            size=8.5, bold=True, color=NAVY)


def table(slide, x, y, w, data, col_w=None, font_size=8.5, head_size=8.5,
          row_h=0.24, head_h=0.26, align=None, fills=None, head_fill=GREEN,
          bold_col0=False):
    """Tableau au style du modèle : en-tête vert, lignes alternées blanc / F5F5F5."""
    nr, nc = len(data), len(data[0])
    gf = slide.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w),
                                Inches(head_h + row_h * (nr - 1)))
    tbl = gf.table
    tbl.first_row = True
    tbl.horz_banding = False
    if col_w:
        tot = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Emu(int(Inches(w) * cw / tot))
    tbl.rows[0].height = Inches(head_h)
    for r in range(1, nr):
        tbl.rows[r].height = Inches(row_h)

    for r, ligne in enumerate(data):
        for c, val in enumerate(ligne):
            cell = tbl.cell(r, c)
            cell.text = ""
            cell.margin_left = cell.margin_right = Inches(0.05)
            cell.margin_top = cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text_frame.word_wrap = True
            p = cell.text_frame.paragraphs[0]
            p.alignment = (align[c] if align else
                           (PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER))
            run = p.add_run()
            run.text = str(val)
            run.font.name = FONT
            entete = r == 0
            run.font.size = Pt(head_size if entete else font_size)
            run.font.bold = entete or (bold_col0 and c == 0)
            run.font.color.rgb = WHITE if entete else BODY
            cell.fill.solid()
            if entete:
                cell.fill.fore_color.rgb = head_fill
            elif fills and (r, c) in fills:
                cell.fill.fore_color.rgb = fills[(r, c)]
                run.font.bold = True
            else:
                cell.fill.fore_color.rgb = WHITE if r % 2 else CARD
    return gf


def encart(slide, x, y, w, h, titre, lignes, ton="green", size=8):
    couleur, fond = {"red": (RED, RED_PALE), "amber": (AMBER, AMBER_PALE),
                     "green": (GREEN, GREEN_SOFT), "navy": (NAVY, CARD)}[ton]
    rect(slide, x, y, w, h, fill=fond)
    rect(slide, x, y, 0.05, h, fill=couleur)
    yy = y + 0.08
    if titre:
        txt(slide, x + 0.16, yy, w - 0.28, 0.22, titre, size=9, bold=True, color=couleur)
        yy += 0.26
    txt(slide, x + 0.16, yy, w - 0.28, max(0.2, y + h - yy - 0.06),
        [f"•  {l}" for l in lignes], size=size, color=BODY, spacing=3.5)


def note(slide, y, texte, size=7.5):
    txt(slide, LEFT, y, 9.2, 0.28, texte, size=size, color=MUTED, italic=True)


def _axes(ch, size, maxval=None, minval=None):
    va = ch.value_axis
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = DIVIDER
    va.major_gridlines.format.line.width = Pt(0.5)
    va.tick_labels.font.size = Pt(size)
    va.tick_labels.font.name = FONT
    va.minimum_scale = minval if minval is not None else 0
    if maxval:
        va.maximum_scale = maxval
    ca = ch.category_axis
    ca.tick_labels.font.size = Pt(size)
    ca.tick_labels.font.name = FONT
    ca.format.line.color.rgb = DIVIDER


def barres(slide, x, y, w, h, categories, series, couleurs=None, horizontal=False,
           maxval=None, minval=None, size=7.5, num_fmt='0"%"', legende=False,
           gap=50, labels=True):
    """series : [(nom, [valeurs]), ...] ; couleurs : [RGBColor] ou [[RGBColor par point]]."""
    cd = CategoryChartData()
    cd.categories = categories
    for nom, vals in series:
        cd.add_series(nom, vals)
    kind = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    ch = slide.shapes.add_chart(kind, Inches(x), Inches(y), Inches(w), Inches(h), cd).chart
    ch.has_title = False
    ch.font.name = FONT
    ch.font.size = Pt(size)
    ch.has_legend = legende
    if legende:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(size)
    plot = ch.plots[0]
    plot.gap_width = gap
    plot.has_data_labels = labels
    if labels:
        dl = plot.data_labels
        dl.number_format = num_fmt
        dl.number_format_is_linked = False
        dl.font.size = Pt(size)
        dl.font.bold = True
        dl.font.name = FONT
        dl.font.color.rgb = BODY
        try:
            dl.position = XL_LABEL_POSITION.OUTSIDE_END
        except Exception:
            pass
    if couleurs:
        for si, ser in enumerate(plot.series):
            if isinstance(couleurs[si], RGBColor):
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = couleurs[si]
            else:
                for pi, pt_ in enumerate(ser.points):
                    pt_.format.fill.solid()
                    pt_.format.fill.fore_color.rgb = couleurs[si][pi]
    _axes(ch, size, maxval, minval)
    return ch


def anneau(slide, x, y, w, h, categories, valeurs, couleurs, size=7.5, labels=True):
    cd = CategoryChartData()
    cd.categories = categories
    cd.add_series("s", valeurs)
    ch = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y),
                                Inches(w), Inches(h), cd).chart
    ch.has_title = False
    ch.font.name = FONT
    ch.font.size = Pt(size)
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False
    ch.legend.font.size = Pt(size)
    plot = ch.plots[0]
    plot.has_data_labels = labels
    if labels:
        dl = plot.data_labels
        dl.number_format = '0.0%'
        dl.number_format_is_linked = False
        dl.show_percentage = True
        dl.show_value = False
        dl.font.size = Pt(size)
        dl.font.bold = True
        dl.font.color.rgb = WHITE
        dl.font.name = FONT
    for pi, pt_ in enumerate(plot.series[0].points):
        pt_.format.fill.solid()
        pt_.format.fill.fore_color.rgb = couleurs[pi % len(couleurs)]
        pt_.format.line.color.rgb = WHITE
        pt_.format.line.width = Pt(1)
    return ch


def fr(n):
    return f"{n:,}".replace(",", " ")


def pc(n, d, dec=1):
    return round(n / d * 100, dec) if d else 0.0


def couleur_taux(v, bas=80, haut=99):
    if v is None:
        return MUTED
    return RED if v < bas else (AMBER if v < haut else GREEN_OK)
